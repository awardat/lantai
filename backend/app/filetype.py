"""文件类型分类与解析：五类文件（文字文档 / Office / 文字 PDF / 图片 / 图片 PDF·OCR）。

分类规则：
- txt / md                     → text      文字文档
- docx                         → office    Office 文档
- pdf（有文本层，≥阈值字符）    → pdf_text  文字 PDF
- pdf（文本层缺失/过少）        → pdf_image 图片 PDF（扫描件，走 OCR）
- png/jpg/jpeg/webp/bmp/gif    → image     图片（视觉模型描述）
"""
from __future__ import annotations

import io
import mimetypes
import re
from pathlib import Path
from typing import Optional

from . import config
from .chunker import clean_ocr_spacing  # 0.1.38 CH-073：OCR 文本层字间空格清理（预览/入库一致）

_SAFE_RE = re.compile(r'[\\/:*?"<>|\x00-\x1f]')


def sanitize_filename(name: str) -> str:
    """净化文件名：去除路径分隔符与控制字符，防止路径穿越。"""
    name = Path(name).name  # 只取最后一段
    name = _SAFE_RE.sub("_", name).strip().strip(".")
    return name or "unnamed"


def classify_ext(ext: str) -> str:
    """扩展名 → 分类（不涉及 PDF 文本层判定）。"""
    ext = ext.lower()
    if ext in (".txt", ".md"):
        return "text"
    if ext in (".docx", ".doc", ".wps", ".xls", ".xlsx", ".ppt", ".pptx"):
        return "office"
    if ext == ".pdf":
        return "pdf_text"  # 占位，解析时再判定是否 pdf_image
    if ext in config.IMAGE_EXTS:
        return "image"
    return "unknown"


def read_text_file(path: Path) -> str:
    """读取 txt/md（显式 UTF-8，容错 GBK 回落）。"""
    raw = path.read_bytes()
    for enc in ("utf-8", "gb18030"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _table_to_nl(rows: list[list[str]], table_name: str) -> list[str]:
    """表格行 → 自然语言句（0.1.43，CH-084，投研报告 V3 方案借鉴）。

    Embedding 对数字/结构化数据区分度差（"15ms"与"20ms"向量几乎重叠），
    表格整块或切碎检索都差；把"表头+行值"转成一句自然语言、每行一句，
    使向量与 BM25 均能精准命中。纯规则模板、零 LLM 依赖（防幻觉）。

    - 首行视为表头；后续每行 → "{表名}：{表头1}为{值1}，{表头2}为{值2}，…"
    - 值非空才生成；表头缺失用"第 N 列"兜底
    - 只有表头（无数据行）→ 原样输出表头
    """
    if not rows:
        return []
    nc = max(len(r) for r in rows)
    header = [str(x).strip() if x is not None else "" for x in (rows[0] if rows else [])]
    out: list[str] = []
    for r in rows[1:]:
        pairs: list[str] = []
        for i in range(nc):
            v = r[i] if i < len(r) else None
            if v is None or str(v).strip() == "":
                continue
            h = header[i] if i < len(header) and header[i] else f"第{i + 1}列"
            pairs.append(f"{h}为{str(v).strip()}")
        if pairs:
            out.append(f"{table_name}：" + "，".join(pairs))
    if not out and header:
        out.append(" | ".join(h for h in header if h))
    return out


def parse_docx(path: Path) -> str:
    """docx → 段落 + 表格文本（表格行转自然语言，CH-084）。"""
    from docx import Document  # python-docx

    doc = Document(str(path))
    parts: list[str] = []
    # 按文档顺序提取段落与表格
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    tbl_idx = 0
    for block in doc.element.body.iterchildren():
        if block.tag.endswith("}p"):
            para = Paragraph(block, doc)
            if para.text.strip():
                parts.append(para.text.strip())
        elif block.tag.endswith("}tbl"):
            tbl_idx += 1
            table = Table(block, doc)
            rows = [[c.text.strip() for c in row.cells] for row in table.rows]
            parts.extend(_table_to_nl(rows, f"表格{tbl_idx}"))
    return "\n".join(parts)


# ---------- 0.1.30：常见办公文档（doc/wps/xls/xlsx/ppt/pptx） ----------

# .doc/.wps 文本提取的最小可读段长度（连续可打印字符数低于此视为噪声丢弃）
_OLE_MIN_RUN = 4


def _ole_utf16le_runs(data: bytes) -> list[str]:
    """从字节流中提取连续可读文本段（UTF-16LE 解码 + 可打印过滤，演示级）。

    OLE2 二进制文档（.doc/.wps）的 WordDocument 流以 UTF-16LE 存储文本，
    简化提取：按 2 字节步长解码，保留连续可打印字符运行。
    """
    runs: list[str] = []
    cur: list[str] = []
    for i in range(0, len(data) - 1, 2):
        cp = data[i] | (data[i + 1] << 8)
        if cp in (0, 0x0D, 0x0A, 0x09) or (0x20 <= cp < 0xFFFF):
            ch = chr(cp)
            if ch.isprintable() or ch in "\r\n\t":
                cur.append(ch)
                continue
        if len(cur) >= _OLE_MIN_RUN:
            runs.append("".join(cur).strip())
        cur = []
    if len(cur) >= _OLE_MIN_RUN:
        runs.append("".join(cur).strip())
    return runs


def parse_doc(path: Path) -> str:
    """.doc（Word 97-2003 OLE2）→ 文本（演示级：WordDocument 流 UTF-16LE 段提取）。"""
    import olefile

    try:
        with olefile.OleFileIO(str(path)) as ole:
            if not ole.exists("WordDocument"):
                return ""
            data = ole.openstream("WordDocument").read()
    except Exception:
        # 非 OLE2 文件（损坏/伪装扩展名）→ 返回空，由上层提示"未能提取文本"
        return ""
    runs = _ole_utf16le_runs(data)
    # 可读性兜底：提取结果过差（如仅噪声）返回空串，由上层提示
    text = "\n".join(runs).strip()
    return text if text_readability(text) >= 0.15 else ""


def parse_wps(path: Path) -> str:
    """.wps（WPS 文字，OLE2 与 Word 兼容）→ 文本（同 .doc 提取）。"""
    return parse_doc(path)


def _format_cell(v, fmt: str = "") -> str:
    """单元格值 → 展示文本（0.1.43 CH-087：格式识别，避免"日期序数/裸百分比"进 NL 句）。

    - datetime/date/time → ISO 文本（"2025-08-27"）
    - 数字 + 百分比格式（如 "0.00%"） → "35.25%"
    - 整型化 float（15.0 → "15"）；其余原样字符串
    """
    import datetime as _dt

    if v is None:
        return ""
    if isinstance(v, bool):
        return "是" if v else "否"
    if isinstance(v, _dt.datetime):
        # 时间全零（Excel 常见）→ 只留日期，避免 "2025-08-27 00:00:00" 噪声
        if v.hour == 0 and v.minute == 0 and v.second == 0 and v.microsecond == 0:
            return v.strftime("%Y-%m-%d")
        return v.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(v, _dt.date):
        return v.strftime("%Y-%m-%d")
    if isinstance(v, _dt.time):
        return v.strftime("%H:%M:%S")
    if isinstance(v, (int, float)) and not isinstance(v, bool):
        if "%" in (fmt or ""):
            return f"{round(float(v) * 100, 6):g}%"
        return str(int(v)) if float(v).is_integer() else str(v)
    return str(v).strip()


def parse_xls(path: Path) -> str:
    """.xls（Excel 97-2003）→ 单元格文本（表格行转自然语言，CH-084；日期转文本 CH-087）。"""
    import xlrd

    try:
        book = xlrd.open_workbook(str(path))
    except Exception:
        return ""  # 非 xls（损坏/伪装扩展名）→ 空，由上层提示
    parts: list[str] = []
    for sheet in book.sheets():
        rows: list[list[str]] = []
        for row in range(sheet.nrows):
            vals: list[str] = []
            for col in range(sheet.ncols):
                cell = sheet.cell(row, col)
                v = cell.value
                if v in (None, ""):
                    vals.append("")
                    continue
                # L2/CH-087：Excel 日期存为 float 序数（如 45200），NL 句内伪语义更"可信"——
                # 按日期类型转换（xlrd 2.x 无 formatting_info，仅能识别日期 ctype）
                if cell.ctype == xlrd.XL_CELL_DATE:
                    d = xlrd.xldate_as_datetime(v, book.datemode)
                    vals.append(d.strftime("%Y-%m-%d"))
                else:
                    vals.append(_format_cell(v))
            if any(x.strip() for x in vals):
                rows.append(vals)
        parts.extend(_table_to_nl(rows, sheet.name or "表格"))
    return "\n".join(parts)


def parse_xlsx(path: Path) -> str:
    """.xlsx → 单元格文本（表格行转自然语言，CH-084；日期/百分比按格式展示 CH-087）。"""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return ""  # 非 xlsx（损坏/伪装扩展名）→ 空，由上层提示
    parts: list[str] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        for row in ws.iter_rows():
            vals = [_format_cell(c.value, c.number_format or "") for c in row]
            if any(x for x in vals):
                rows.append(vals)
        parts.extend(_table_to_nl(rows, ws.title or "表格"))
    wb.close()
    return "\n".join(parts)


def _shape_text(shape) -> str:
    """递归提取 shape 文本（含文本框/表格/分组）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[str] = []
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            out.append(_shape_text(sub))
        return "\n".join(out)
    if st == MSO_SHAPE_TYPE.TABLE:
        # 0.1.43（CH-084）：表格行转自然语言
        rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
        return "\n".join(_table_to_nl(rows, "表格"))
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                out.append(t)
    return "\n".join(out)


def parse_pptx(path: Path) -> str:
    """.pptx → 幻灯片文本框/表格文本。"""
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception:
        return ""  # 非 pptx（损坏/伪装扩展名）→ 空，由上层提示
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            t = _shape_text(shape).strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


def parse_office(path: Path, ext: str) -> str:
    """office 类统一解析入口（按扩展名分发，0.1.30）。"""
    ext = ext.lower()
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".doc":
        return parse_doc(path)
    if ext == ".wps":
        return parse_wps(path)
    if ext == ".xls":
        return parse_xls(path)
    if ext == ".xlsx":
        return parse_xlsx(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext == ".ppt":
        return ""  # .ppt 为 OLE2 二进制演示文稿，暂无可靠纯 Python 提取，返回空由上层提示
    return ""


_CHAR_CODE_TOKEN_RE = re.compile(r"/[A-Za-z0-9]{1,8}")
_WORD_RE = re.compile(r"[A-Za-z]{3,}")


def text_readability(text: str) -> float:
    """0~1 文本可读性：CJK 占比 + 完整英文单词占比，扣减字符码 token（如 /G21/G22）。

    用途：区分"真实文本"与"编码不可映射的字符码伪文本"（内嵌字体缺 ToUnicode 的 PDF）。
    """
    if not text:
        return 0.0
    total = len(text)
    cjk = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff" or "\u3400" <= ch <= "\u4dbf")
    word_chars = sum(len(w) for w in _WORD_RE.findall(text))
    code_tokens = len(_CHAR_CODE_TOKEN_RE.findall(text))
    score = (cjk * 1.5 + word_chars) / total + 0.5 - code_tokens * 4 / total
    return max(0.0, min(1.0, score))


# 可读性阈值（0.1.16 样本验证）：GBT 43052 伪文本（字符码）得分 0.0；
# 正常中文/英文文本得分约 1.0；0.4 留出宽松余量，兼顾英文文档与少量噪声。
READABILITY_THRESHOLD = 0.4


def pdf_is_pseudo_text(path: Path) -> bool:
    """检测"文本层存在但编码不可映射"（pypdf 能提取出字符码伪文本，可读性低）。"""
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        full = "\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        return False
    return len(full) >= config.PDF_TEXT_MIN_CHARS and text_readability(full) < READABILITY_THRESHOLD


def _pypdf_page_text(path: Path, page_index: int, reader=None) -> str:
    """pypdf 单页文本；reader 可复用（0.1.38 审核 L1：避免 N 页 N 次全文解析）。"""
    if reader is None:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
    if page_index >= len(reader.pages):
        return ""
    try:
        return reader.pages[page_index].extract_text() or ""
    except Exception:
        return ""


def pdf_text_layers(path: Path) -> list[tuple[str, bool]]:
    """逐页提取文本（R117①：pdfminer layout 按位置取行 → 几何排序还原阅读顺序）。

    返回 [(页面文本, 是否文本页), ...]：
    - 页面文本：按阅读顺序（栏内 y 降序、x 升序）拼接的行；
    - 是否文本页：该页有效文本 ≥ PDF_TEXT_MIN_CHARS 字符且可读性达标（0.1.16 双引擎兜底：
      pdfminer 空/不可读时回退 pypdf，仍不可读（如编码不可映射的伪文本）→ 判图片页走 OCR）。
    处理：跨页重复行（页眉/页脚模式）剔除、页顶/页底纯数字页码碎片剔除、基础双栏检测。
    """
    from collections import Counter

    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextLineHorizontal, LTTextContainer

    # 0.1.38（CH-072）：pdfminer 打开失败（startxref/xref 非标，如"国办发2014 6号"
    # 通知 PDF 抛 PDFSyntaxError "No /Root object!"——阅读器/pypdf 均可打开）→
    # 回退 pypdf 逐页提取，不再让预览/解析 500
    try:
        page_layouts = list(extract_pages(str(path)))
    except Exception:  # noqa: BLE001
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        out: list[tuple[str, bool]] = []
        for idx in range(len(reader.pages)):
            # 0.1.38（CH-073）：回退分支同样清理 OCR 字间空格（预览/入库一致）；
            # 复用 reader（审核 L1：避免每页重开全文解析）
            t = clean_ocr_spacing(_pypdf_page_text(path, idx, reader)).strip()
            if len(t) >= config.PDF_TEXT_MIN_CHARS and text_readability(t) >= READABILITY_THRESHOLD:
                out.append((t, True))
            else:
                out.append(("", False))
        return out

    page_lines: list[list[tuple[str, float, float, float, float]]] = []
    all_lines: list[str] = []
    heights: list[float] = []
    for page_layout in page_layouts:
        lines: list[tuple[str, float, float, float, float]] = []
        for el in page_layout:
            if not isinstance(el, LTTextContainer):
                continue
            for ln in el:
                if isinstance(ln, LTTextLineHorizontal):
                    t = ln.get_text().strip()
                    if t:
                        lines.append((t, ln.x0, ln.x1, ln.y0, ln.y1))
        page_lines.append(lines)
        all_lines.extend(t for t, *_ in lines)
        heights.append(float(page_layout.height))
    # 跨页重复行（≥3 页出现 → 页眉/页脚等装饰性内容）
    repeated = {t for t, c in Counter(all_lines).items() if c >= 3}

    pages_out: list[tuple[str, bool]] = []
    pypdf_reader = None  # 审核 L1：主路径 fallback 复用单个 reader，避免每页全文解析
    for idx, (lines, height) in enumerate(zip(page_lines, heights)):
        ordered = _layout_order(lines, height, repeated)
        # 0.1.38（CH-073 方案 A）：OCR 文本层字间空格/换行碎片清理（预览与入库一致）
        text = clean_ocr_spacing("\n".join(ordered)).strip()
        if len(text) >= config.PDF_TEXT_MIN_CHARS and text_readability(text) >= READABILITY_THRESHOLD:
            pages_out.append((text, True))
            continue
        # 双引擎兜底：pdfminer 空/不可读 → 尝试 pypdf（0.1.16）
        if pypdf_reader is None:
            from pypdf import PdfReader

            pypdf_reader = PdfReader(str(path))
        fallback = clean_ocr_spacing(_pypdf_page_text(path, idx, pypdf_reader)).strip()
        if len(fallback) >= config.PDF_TEXT_MIN_CHARS and text_readability(fallback) >= READABILITY_THRESHOLD:
            pages_out.append((fallback, True))
        else:
            pages_out.append((text, False))
    return pages_out


def _layout_order(
    lines: list[tuple[str, float, float, float, float]],
    height: float,
    repeated: set[str],
) -> list[str]:
    """按阅读顺序排列文本行：过滤装饰 → 双栏检测 → 栏内几何排序。"""
    kept: list[tuple[str, float, float, float, float]] = []
    for t, x0, x1, y0, y1 in lines:
        s = t.strip()
        if not s:
            continue
        if s in repeated:  # 跨页重复（页眉/页脚）
            continue
        # 页码碎片：纯数字且位于页顶（上 12%）/页底（下 15%）区域
        if re.fullmatch(r"[\d\s]{1,6}", s) and (y1 > height * 0.88 or y0 < height * 0.15):
            continue
        kept.append((t, x0, x1, y0, y1))
    if not kept:
        return []

    # 基础双栏检测：行中点 x 的最大空隙
    mids = sorted((x0 + x1) / 2 for _, x0, x1, _, _ in kept)
    page_w = max(x1 for _, _, x1, _, _ in kept) - min(x0 for _, x0, _, _, _ in kept)
    gap, split = 0.0, None
    for a, b in zip(mids, mids[1:]):
        if b - a > gap:
            gap, split = b - a, (a + b) / 2
    two_column = gap > page_w * 0.12 and len(kept) >= 16

    def by_col(col: str) -> list[str]:
        rows = [k for k in kept if (k[1] + k[2]) / 2 < split] if col == "left" else [k for k in kept if (k[1] + k[2]) / 2 >= split]
        return [t for t, *_ in sorted(rows, key=lambda k: (-k[4], k[1]))]  # y 降序、x 升序

    if two_column:
        return by_col("left") + by_col("right")
    return [t for t, *_ in sorted(kept, key=lambda k: (-k[4], k[1]))]


def pdf_extract_page_images(path: Path) -> list[tuple[int, bytes, str]]:
    """提取 PDF 各页内嵌图片（扫描件通常为整页图）。返回 [(页码, 字节, media_type)]。

    依赖 Pillow（pypdf 图片提取要求 `pip install pypdf[image]`）。
    """
    from pypdf import PdfReader

    try:
        from PIL import Image  # noqa: F401  确保 pillow 可用
    except ImportError as exc:
        raise RuntimeError("PDF 图片提取需要 Pillow，请执行：pip install pillow（requirements 已包含）。") from exc

    reader = PdfReader(str(path))
    out: list[tuple[int, bytes, str]] = []
    for i, page in enumerate(reader.pages):
        try:
            for img in page.images:
                data = img.data
                mime = mimetypes.guess_type(img.name or "x.png")[0] or "image/png"
                out.append((i + 1, data, mime))
        except Exception:
            continue
    return out


def mime_of(ext: str) -> str:
    return mimetypes.guess_type("x" + ext)[0] or "application/octet-stream"


# 视觉供应商通用支持格式（MiMo/通义/OpenAI 等均支持；超出即 400 拒绝）
_VISION_OK_FORMATS = {"JPEG", "PNG", "BMP", "WEBP", "GIF"}


def pdf_render_page_images(path: Path, zoom: float = 2.0) -> list[tuple[int, bytes, str]]:
    """pymupdf 整页渲染为 PNG（0.1.38，CH-070）：用于无内嵌位图的矢量描摹/混合 PDF，
    作为 OCR 通道兜底（R110 整页渲染评估项落地）。返回 [(页码, PNG 字节, image/png)]。

    依赖 pymupdf（见 requirements；Python 3.14 用 cp310-abi3 wheel）。
    """
    import pymupdf  # 新版 API（fitz 已弃用）

    doc = pymupdf.open(str(path))
    out: list[tuple[int, bytes, str]] = []
    try:
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom))
            out.append((i, pix.tobytes("png"), "image/png"))
    finally:
        doc.close()
    return out


def normalize_image_for_vision(data: bytes, mime: str) -> tuple[bytes, str]:
    """视觉模型图片标准化（0.1.37，CH-066）：Pillow 探测图片实际格式，不在供应商通用
    支持范围（jpeg/png/bmp/webp/gif）的格式（如 PDF 内嵌 TIFF/JPEG2000/CCITT）统一
    转码为 JPEG，消除"仅支持 jpg、bmp、webp、gif、png 格式"类 400。解码失败原样返回
    （交由上游报错，不吞异常）。
    """
    from PIL import Image

    try:
        with Image.open(io.BytesIO(data)) as im:
            if (im.format or "").upper() in _VISION_OK_FORMATS:
                return data, mime  # 已支持，原样发送（避免无谓转码与体积膨胀）
            buf = io.BytesIO()
            im.convert("RGB").save(buf, format="JPEG", quality=90)
            return buf.getvalue(), "image/jpeg"
    except Exception:  # noqa: BLE001 无法解码（损坏/未知）→ 原样发送，由上游返回具体错误
        return data, mime
